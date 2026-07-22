import os
import fitz # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
import speech_recognition as sr
from pydub import AudioSegment
import imageio_ffmpeg
import uuid

# Configure pydub to use the ffmpeg binary from imageio_ffmpeg
AudioSegment.converter = imageio_ffmpeg.get_ffmpeg_exe()

# Maximum number of characters to extract to prevent AI context overflow
MAX_CHARS = 12000

def parse_document(file_path: str, mime_type: str) -> str:
    """Extracts text from a document file based on its MIME type."""
    ext = os.path.splitext(file_path)[1].lower()
    
    text = ""
    try:
        if ext == ".pdf" or "pdf" in mime_type:
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text() + "\n"
                if len(text) > MAX_CHARS:
                    break
                    
        elif ext in [".docx", ".doc"] or "word" in mime_type:
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            
        elif ext in [".xlsx", ".xls", ".csv"] or "excel" in mime_type or "spreadsheet" in mime_type:
            if ext == ".csv":
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            text = df.head(100).to_string() # Read first 100 rows
            
        elif ext in [".pptx", ".ppt"] or "powerpoint" in mime_type or "presentation" in mime_type:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) > MAX_CHARS:
                    break
        else:
            return "Unsupported document format for text extraction."
            
    except Exception as e:
        print(f"Error parsing document: {e}")
        return f"Error reading file: {e}"
        
    if not text.strip():
        return "The document appears to be empty or contains no readable text (it might be scanned images)."
        
    return text[:MAX_CHARS]


def transcribe_audio_video(file_path: str) -> str:
    """Extracts audio and transcribes 100% of spoken speech for media up to 3 hours in duration."""
    temp_wav = f"temp_{uuid.uuid4().hex}.wav"
    temp_mp3 = f"temp_{uuid.uuid4().hex}.mp3"
    
    try:
        # Extract audio track - use 24kHz 64kbps for clearer speech fidelity
        audio = AudioSegment.from_file(file_path)
        
        # Limit max duration to 3 hours (10,800,000 ms = 180 minutes)
        max_3hr_ms = 3 * 3600 * 1000
        if len(audio) > max_3hr_ms:
            audio = audio[:max_3hr_ms]
            
        # 24kHz mono, 64kbps — high enough quality for accurate speech recognition
        audio_mono = audio.set_frame_rate(24000).set_channels(1)
        audio_mono.export(temp_mp3, format="mp3", bitrate="64k")
        
        # 1. Try Gemini Flash for exact verbatim transcription (temperature=0)
        api_key = os.environ.get("GEMINI_API_KEY")
        if api_key and os.path.exists(temp_mp3):
            try:
                import base64
                import requests
                
                # If audio is very long (>30 min), split into 30-minute parts for Gemini REST payload
                gemini_transcripts = []
                part_duration_ms = 30 * 60 * 1000  # 30 mins
                total_len = len(audio_mono)
                
                for p_start in range(0, total_len, part_duration_ms):
                    part_audio = audio_mono[p_start : min(p_start + part_duration_ms, total_len)]
                    temp_part = f"temp_part_{uuid.uuid4().hex}.mp3"
                    part_audio.export(temp_part, format="mp3", bitrate="64k")
                    
                    try:
                        with open(temp_part, "rb") as f:
                            b64 = base64.b64encode(f.read()).decode("utf-8")
                        
                        # Use gemini-2.0-flash for best accuracy, fallback to 1.5-flash
                        for model in ["gemini-2.0-flash", "gemini-1.5-flash"]:
                            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
                            payload = {
                                "generationConfig": {
                                    "temperature": 0.0,
                                    "topP": 0.05,
                                    "topK": 1
                                },
                                "contents": [{
                                    "parts": [
                                        {
                                            "text": (
                                                "You are a VERBATIM speech-to-text transcription engine with zero tolerance for errors.\n"
                                                "TASK: Transcribe EVERY SINGLE WORD spoken in this audio EXACTLY as it was said.\n"
                                                "CRITICAL RULES:\n"
                                                "1. Write ONLY what is actually spoken — never substitute, paraphrase, or guess a different word.\n"
                                                "2. If a speaker says 'apple', you MUST write 'apple' — NOT 'banana', NOT 'fruit', NOT any other word.\n"
                                                "3. Detect the spoken language automatically (Khmer, English, Chinese, Vietnamese, Thai, Korean, Japanese, French, Spanish, German, Russian, Arabic, Hindi, etc.) and transcribe in the ORIGINAL language script.\n"
                                                "4. Do NOT add titles, headers, timestamps, descriptions, or any commentary — output the raw spoken words ONLY.\n"
                                                "5. Preserve all proper nouns, names, numbers, and punctuation exactly as spoken.\n"
                                                "6. If a section is inaudible, write [inaudible] — do NOT invent words."
                                            )
                                        },
                                        {"inline_data": {"mime_type": "audio/mp3", "data": b64}}
                                    ]
                                }]
                            }
                            try:
                                res = requests.post(url, json=payload, timeout=120).json()
                                if "candidates" in res and res["candidates"]:
                                    p_text = res["candidates"][0]["content"]["parts"][0]["text"].strip()
                                    if p_text and len(p_text) > 3:
                                        gemini_transcripts.append(p_text)
                                        break  # success, stop trying other models
                            except Exception as model_e:
                                print(f"Gemini {model} transcription error: {model_e}")
                                continue
                    finally:
                        if os.path.exists(temp_part):
                            os.remove(temp_part)
                            
                full_gemini_text = " ".join(gemini_transcripts).strip()
                if full_gemini_text and len(full_gemini_text) > 3:
                    return full_gemini_text
            except Exception as e:
                print(f"Gemini Multi-part Audio Transcription Error: {e}")

        # 2. Multi-language Chunked Google SpeechRecognition Fallback (Supports All Spoken Languages)
        audio_mono.export(temp_wav, format="wav")
        recognizer = sr.Recognizer()
        
        chunk_length_ms = 15000  # 15-second chunks
        total_length_ms = len(audio_mono)
        full_transcript = []
        
        # Supported global languages list
        global_langs = [
            "km-KH", "en-US", "zh-CN", "vi-VN", "th-TH", "ko-KR", "ja-JP", 
            "fr-FR", "es-ES", "de-DE", "ru-RU", "id-ID", "hi-IN", "ar-SA", 
            "pt-PT", "it-IT", "tr-TR", "ms-MY", "my-MM", "fil-PH"
        ]
        
        # Determine primary language by testing first 15 seconds
        first_chunk = audio_mono[:min(15000, total_length_ms)]
        temp_first = f"temp_first_{uuid.uuid4().hex}.wav"
        first_chunk.export(temp_first, format="wav")
        
        detected_lang = "km-KH"
        with sr.AudioFile(temp_first) as source:
            audio_data = recognizer.record(source)
            for lang in global_langs:
                try:
                    t = recognizer.recognize_google(audio_data, language=lang)
                    if t and len(t.strip()) > 2:
                        detected_lang = lang
                        break
                except:
                    continue
        if os.path.exists(temp_first):
            os.remove(temp_first)

        # Transcribe every 15-second chunk sequentially with language fallback
        with sr.AudioFile(temp_wav) as source:
            for start_ms in range(0, total_length_ms, chunk_length_ms):
                try:
                    duration_sec = min(15, (total_length_ms - start_ms) / 1000.0)
                    if duration_sec <= 0.5:
                        break
                    audio_data = recognizer.record(source, duration=duration_sec)
                    
                    chunk_success = False
                    # Try detected_lang first, then fallback to global languages
                    for l_try in [detected_lang] + [l for l in global_langs if l != detected_lang]:
                        try:
                            chunk_text = recognizer.recognize_google(audio_data, language=l_try)
                            if chunk_text:
                                full_transcript.append(chunk_text)
                                chunk_success = True
                                break
                        except sr.UnknownValueError:
                            continue
                        except Exception:
                            continue
                except Exception as e_chunk:
                    print(f"Chunk error at {start_ms}ms: {e_chunk}")
                    
        result_text = " ".join(full_transcript).strip()
        if result_text:
            return result_text
            
        return "Could not recognize speech in media."
    except Exception as e:
        print(f"Error transcribing media: {e}")
        return f"Error extracting speech: {e}"
    finally:
        if os.path.exists(temp_wav): os.remove(temp_wav)
        if os.path.exists(temp_mp3): os.remove(temp_mp3)


def transcribe_with_timestamps(file_path: str) -> list:
    """
    Returns a list of (start_sec, end_sec, text) tuples from Gemini with precise timestamps.
    Used for segment-level timed dubbing so translated voice plays exactly when original speaker speaks.
    Falls back to empty list if Gemini timestamps are unavailable.
    """
    import base64
    import requests
    import re

    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        return []

    temp_mp3 = f"temp_ts_{uuid.uuid4().hex}.mp3"
    try:
        audio = AudioSegment.from_file(file_path)
        # Limit to 30 min for Gemini payload size
        max_ms = 30 * 60 * 1000
        if len(audio) > max_ms:
            audio = audio[:max_ms]
        audio_mono = audio.set_frame_rate(24000).set_channels(1)
        audio_mono.export(temp_mp3, format="mp3", bitrate="64k")

        with open(temp_mp3, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("utf-8")

        # Try gemini-2.0-flash first, then 1.5-flash
        for model in ["gemini-2.0-flash", "gemini-1.5-flash"]:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
            payload = {
                "generationConfig": {"temperature": 0.0, "topP": 0.05, "topK": 1},
                "contents": [{
                    "parts": [
                        {
                            "text": (
                                "Transcribe this audio with EXACT start and end timestamps for every spoken sentence or phrase.\n"
                                "Output ONLY in this strict format, one line per segment:\n"
                                "[HH:MM:SS.mmm --> HH:MM:SS.mmm] spoken text here\n\n"
                                "Rules:\n"
                                "- Include ALL spoken words verbatim in original language.\n"
                                "- Timestamps must match exactly when speech starts and ends.\n"
                                "- Do NOT add any intro, outro, or explanation text.\n"
                                "- If multiple speakers, still use same format."
                            )
                        },
                        {"inline_data": {"mime_type": "audio/mp3", "data": b64}}
                    ]
                }]
            }
            try:
                res = requests.post(url, json=payload, timeout=120).json()
                if "candidates" not in res or not res["candidates"]:
                    continue
                raw_text = res["candidates"][0]["content"]["parts"][0]["text"].strip()

                # Parse [HH:MM:SS.mmm --> HH:MM:SS.mmm] text
                pattern = re.compile(
                    r'\[(\d+):(\d+):([\d.]+)\s*-->\s*(\d+):(\d+):([\d.]+)\]\s*(.+)'
                )
                segments = []
                for m in pattern.finditer(raw_text):
                    sh, sm, ss, eh, em, es, txt = m.groups()
                    start = int(sh) * 3600 + int(sm) * 60 + float(ss)
                    end   = int(eh) * 3600 + int(em) * 60 + float(es)
                    txt = txt.strip()
                    if txt and end > start:
                        segments.append((start, end, txt))

                if segments:
                    print(f"[Timestamps] Got {len(segments)} segments from {model}")
                    return segments
            except Exception as model_e:
                print(f"Timestamp transcription error ({model}): {model_e}")
                continue

    except Exception as e:
        print(f"transcribe_with_timestamps outer error: {e}")
    finally:
        if os.path.exists(temp_mp3):
            os.remove(temp_mp3)

    return []
